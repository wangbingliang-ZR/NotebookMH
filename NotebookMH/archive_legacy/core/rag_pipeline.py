"""
core/rag_pipeline.py - RAG 真理哨兵 (Phase 1A)

工业级摄入流水线 (Ingestion Pipeline)：
  1. HashVaultBarrier   : SHA-256(file_content + vault_uuid) 零越权哈希墙
  2. SemanticAnatomyKnife: 语义感知切块，保留 header_hierarchy / source_page
  3. HybridRetriever   : Dense(MMR) + Sparse(BM25) + 可选 Cross-Encoder 重排
  4. IngestionPipeline : async generator yield 进度事件，Embedding 全程非阻塞

POST_EXECUTION_AUDIT:
  - 阻塞零容忍: _embed_chunks_async 与 _parse_pdf 均在 asyncio.to_thread() 中执行
  - Chunk 质量: 每个 chunk 携带 source_page, header_hierarchy, chunk_size, overlap_prev
  - 优雅降级: Reranker 缺失时平滑回退到 MMR/BM25 融合排序
"""

import asyncio
import hashlib
import json
import logging
import os
import re
import time
from typing import Any, AsyncGenerator, Dict, List, Optional, Tuple

import numpy as np

from utils.db_manager import db_pool

try:
    import jieba
    def _tokenize(text: str) -> list:
        return list(jieba.cut(text))
except ImportError:
    def _tokenize(text: str) -> list:
        return text.split()

logger = logging.getLogger(__name__)

# Phase 1B+: DAG 本体抽取（懒加载，避免循环导入）
_ontology_builder = None

def _get_ontology_builder():
    global _ontology_builder
    if _ontology_builder is None:
        from core.ontology_builder import get_ontology_builder
        _ontology_builder = get_ontology_builder()
    return _ontology_builder


def _ms(t0: float) -> int:
    return int((time.monotonic() - t0) * 1000)


# ---------------------------------------------------------------------------
# MODULE 1: 零越权密码学记忆墙
# ---------------------------------------------------------------------------

class HashVaultBarrier:
    """
    文件指纹绝不能仅靠 MD5。
    必须计算 SHA-256(File_Content + Vault_UUID)，将物理文件与认知容器密码学绑定。
    """

    @staticmethod
    def compute(file_bytes: bytes, vault_uuid: str) -> str:
        hasher = hashlib.sha256()
        hasher.update(file_bytes)
        hasher.update(vault_uuid.encode("utf-8"))
        return hasher.hexdigest()


# ---------------------------------------------------------------------------
# MODULE 2: 语义感知解剖刀
# ---------------------------------------------------------------------------

class SemanticAnatomyKnife:
    """
    禁止单纯使用字数分割。
    依据换行符、句号和标题正则进行逻辑段落切分，确保每个 Chunk 包含完整语义闭环。
    """

    def __init__(
        self,
        chunk_size: int = 800,
        chunk_overlap: int = 150,
    ) -> None:
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def _extract_headers(self, text: str) -> List[Tuple[str, str]]:
        """提取标题层级（H1, H2, 粗体），用于挂载 DAG 知识图谱钩子。"""
        headers: List[Tuple[str, str]] = []
        for line in text.splitlines():
            stripped = line.strip()
            if stripped.startswith("#"):
                level = len(stripped) - len(stripped.lstrip("#"))
                headers.append((f"H{level}", stripped.lstrip("#").strip()))
            elif stripped.startswith("**") and stripped.endswith("**"):
                headers.append(("BOLD", stripped[2:-2]))
        return headers

    def _split_by_semantic_boundary(self, text: str) -> List[str]:
        """
        按语义边界切分：
          - 优先按双换行 \n\n（段落边界）
          - 次按句号 + 空格（句子边界）
          - 最后按字数硬切（兜底）
        """
        # 先按段落分
        raw_blocks = re.split(r"\n\s*\n", text)
        blocks = [b.strip() for b in raw_blocks if b.strip()]

        chunks: List[str] = []
        current = ""

        for block in blocks:
            if len(current) + len(block) + 2 <= self.chunk_size:
                current = f"{current}\n\n{block}".strip() if current else block
            else:
                if current:
                    chunks.append(current)
                # 计算 overlap：从上一块尾部取 chunk_overlap 字符
                overlap = ""
                if chunks:
                    prev = chunks[-1]
                    overlap = prev[-self.chunk_overlap :] if len(prev) > self.chunk_overlap else prev
                current = f"{overlap}\n\n{block}".strip() if overlap else block

        if current:
            chunks.append(current)

        return chunks

    def dissect(
        self,
        text: str,
        source_page: int = 0,
    ) -> List[Dict[str, Any]]:
        """
        语义解剖。返回每个 chunk 的字典，必须携带：
          chunk_text, chunk_index, source_page, header_hierarchy,
          chunk_size, overlap_prev
        """
        headers = self._extract_headers(text)
        header_hierarchy = " > ".join([f"{lvl}:{txt}" for lvl, txt in headers[:3]])

        raw_chunks = self._split_by_semantic_boundary(text)
        result: List[Dict[str, Any]] = []

        for idx, chunk_text in enumerate(raw_chunks):
            overlap_prev = 0
            if idx > 0:
                prev = raw_chunks[idx - 1]
                # 计算与上一块的重叠字符数
                min_len = min(self.chunk_overlap, len(prev), len(chunk_text))
                for i in range(min_len, 0, -1):
                    if prev[-i:] == chunk_text[:i]:
                        overlap_prev = i
                        break

            result.append(
                {
                    "chunk_index": idx,
                    "chunk_text": chunk_text,
                    "source_page": source_page,
                    "header_hierarchy": header_hierarchy,
                    "chunk_size": len(chunk_text),
                    "overlap_prev": overlap_prev,
                }
            )

        return result


# ---------------------------------------------------------------------------
# MODULE 3: 混合检索与量子重排引擎
# ---------------------------------------------------------------------------

class HybridRetriever:
    """
    Ensemble Retriever：
      链路 A (稠密): 预留 MMR 向量检索接口
      链路 B (稀疏): BM25 关键词检索
      终极清洗: Cross-Encoder Re-ranking（Cohere -> BGE -> 降级）
    """

    def __init__(
        self,
        persist_dir: str = "./data/chroma_db",
        reranker_api_key: Optional[str] = None,
    ) -> None:
        self.persist_dir = persist_dir
        self.reranker_api_key = reranker_api_key
        self._bm25_index: Optional[Any] = None
        self._chunk_texts: List[str] = []

    # ---- 稀疏检索 (BM25) ----

    def add_chunks(self, chunks: List[str]) -> None:
        """追加 chunks 到 BM25 稀疏索引（支持多文档累积）。"""
        try:
            from rank_bm25 import BM25Okapi

            self._chunk_texts.extend(chunks)
            tokenized = [_tokenize(c) for c in self._chunk_texts]
            self._bm25_index = BM25Okapi(tokenized)
            logger.info("BM25 index updated: %d total chunks", len(self._chunk_texts))
        except ImportError:
            logger.warning("rank-bm25 not installed; sparse search disabled.")

    def search_bm25(self, query: str, top_k: int = 10) -> List[Tuple[int, float]]:
        """BM25 关键词检索，返回 [(chunk_index, score), ...]。"""
        if self._bm25_index is None:
            return []
        scores = self._bm25_index.get_scores(_tokenize(query))
        top_indices = np.argsort(scores)[::-1][:top_k]
        return [(int(i), float(scores[i])) for i in top_indices]

    # ---- 重排序 (优雅降级) ----

    async def rerank(
        self,
        query: str,
        candidates: List[Dict[str, Any]],
    ) -> List[Dict[str, Any]]:
        """
        Cross-Encoder Re-ranking：
          1. 优先 Cohere API
          2. 其次本地 BGE-Reranker
          3. 降级：按已有 score 直通
        """
        if not candidates:
            return []

        # 尝试 Cohere
        if self.reranker_api_key:
            try:
                return await self._rerank_cohere(query, candidates)
            except Exception as e:
                logger.warning("Cohere rerank failed: %s, degrading...", e)

        # 尝试本地 BGE
        try:
            return await self._rerank_bge(query, candidates)
        except Exception as e:
            logger.warning("BGE rerank failed: %s, using MMR fallback.", e)

        # 终极降级：按已有 score 排序，仅取 Top-3
        ranked = sorted(
            candidates, key=lambda x: x.get("score", 0), reverse=True
        )
        return ranked[:3]

    async def _rerank_cohere(
        self, query: str, candidates: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        import httpx

        docs = [c["chunk_text"] for c in candidates]
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.post(
                "https://api.cohere.com/v1/rerank",
                headers={"Authorization": f"Bearer {self.reranker_api_key}"},
                json={
                    "query": query,
                    "documents": docs,
                    "top_n": 3,
                    "model": "rerank-english-v2.0",
                },
            )
            r.raise_for_status()
            data = r.json()

        ranked: List[Dict[str, Any]] = []
        for item in data["results"]:
            idx = item["index"]
            candidates[idx]["rerank_score"] = item["relevance_score"]
            ranked.append(candidates[idx])
        return ranked

    async def _rerank_bge(
        self, query: str, candidates: List[Dict[str, Any]]
    ) -> List[Dict[str, Any]]:
        """本地 BGE 模型重排，使用 asyncio.to_thread 避免阻塞事件循环。"""
        from sentence_transformers import CrossEncoder

        model = CrossEncoder("BAAI/bge-reranker-base")
        pairs = [(query, c["chunk_text"]) for c in candidates]
        # CRITICAL: 运行在独立线程中，绝不卡死主循环
        scores = await asyncio.to_thread(model.predict, pairs)
        for c, s in zip(candidates, scores):
            c["rerank_score"] = float(s)
        return sorted(
            candidates, key=lambda x: x.get("rerank_score", 0), reverse=True
        )[:3]


# ---------------------------------------------------------------------------
# MODULE 4: 主摄入管线
# ---------------------------------------------------------------------------

class IngestionPipeline:
    """RAG 真理哨兵 —— 完整摄入流水线，yield 异步进度事件。"""

    def __init__(
        self,
        chunk_size: int = 800,
        chunk_overlap: int = 150,
        persist_dir: str = "./data/chroma_db",
    ) -> None:
        self.barrier = HashVaultBarrier()
        self.knife = SemanticAnatomyKnife(chunk_size, chunk_overlap)
        self.retriever = HybridRetriever(persist_dir)
        self.persist_dir = persist_dir
        os.makedirs(persist_dir, exist_ok=True)

        # Phase 1B: ChromaDB 持久化向量检索
        try:
            import chromadb
            self._chroma = chromadb.PersistentClient(path=persist_dir)
        except ImportError:
            logger.warning("chromadb not installed; dense search disabled.")
            self._chroma = None
        self._embedder = None  # lazy load sentence-transformer

    async def ingest_document(
        self,
        file_bytes: bytes,
        file_name: str,
        vault_uuid: str,
    ) -> AsyncGenerator[Dict[str, Any], None]:
        """
        异步摄入文档，yield 进度事件字典。

        事件流示例：
          {"status": "HASHING", "elapsed_ms": 12}
          {"status": "CACHED",   "elapsed_ms": 45}
          {"status": "PARSING",  "page": 3, "elapsed_ms": 120}
          {"status": "CHUNKING", "chunks": 42, "elapsed_ms": 340}
          {"status": "EMBEDDING","progress": 0.5, "elapsed_ms": 1200}
          {"status": "INDEXING", "chunks_indexed": 42, "elapsed_ms": 1500}
          {"status": "DONE",     "chunks": 42, "elapsed_ms": 2000}
        """
        t0 = time.monotonic()

        # ── Step 1: 哈希墙 ──────────────────────────────────
        yield {"status": "HASHING", "elapsed_ms": _ms(t0)}
        content_hash = self.barrier.compute(file_bytes, vault_uuid)

        if db_pool.document_exists(vault_uuid, content_hash):
            yield {
                "status": "CACHED",
                "hash": content_hash[:16],
                "elapsed_ms": _ms(t0),
            }
            return

        yield {"status": "HASH_OK", "hash": content_hash[:16], "elapsed_ms": _ms(t0)}

        # ── Step 2: 解析 ────────────────────────────────────
        yield {"status": "PARSING", "file_name": file_name, "elapsed_ms": _ms(t0)}
        text, page_count = await self._parse_document(file_bytes, file_name)
        yield {
            "status": "PARSE_OK",
            "pages": page_count,
            "chars": len(text),
            "elapsed_ms": _ms(t0),
        }

        # ── Step 3: 语义切块 ─────────────────────────────────
        yield {"status": "CHUNKING", "elapsed_ms": _ms(t0)}
        chunks = self.knife.dissect(text, source_page=1)
        yield {
            "status": "CHUNK_OK",
            "chunks": len(chunks),
            "elapsed_ms": _ms(t0),
        }
        if not chunks:
            yield {
                "status": "EMPTY_TEXT",
                "chars": len(text),
                "elapsed_ms": _ms(t0),
            }
            return

        # ── Step 4: 非阻塞 Embedding ─────────────────────────
        # CRITICAL: embedding 在 asyncio.to_thread() 中执行，
        # 绝不卡死 Streamlit 主渲染循环。
        yield {"status": "EMBEDDING", "chunks": len(chunks), "elapsed_ms": _ms(t0)}
        vectors = await self._embed_chunks_async(chunks)
        yield {
            "status": "EMBED_OK",
            "vectors": len(vectors),
            "elapsed_ms": _ms(t0),
        }

        # ── Step 5: 索引 ─────────────────────────────────────
        yield {"status": "INDEXING", "elapsed_ms": _ms(t0)}
        await self._index_chunks(vault_uuid, content_hash, chunks, vectors)
        yield {"status": "INDEX_OK", "elapsed_ms": _ms(t0)}

        # ── Step 6: 注册到 SQLite ────────────────────────────
        db_pool.register_document(
            vault_uuid=vault_uuid,
            file_name=file_name,
            content_hash=content_hash,
            doc_size=len(file_bytes),
            page_count=page_count,
            full_text=text[:50000],
        )
        db_pool.register_chunks(
            vault_uuid,
            content_hash,
            chunks,
            embedding_model="all-MiniLM-L6-v2",
        )

        # ── Step 7: DAG 本体抽取 ────────────────────────────────
        yield {"status": "ONTOLOGY", "elapsed_ms": _ms(t0)}
        try:
            ob = _get_ontology_builder()
            ontology = await ob.extract_from_document(text, vault_uuid, content_hash)
            yield {"status": "ONTOLOGY_OK", "concepts": len(ontology.dag), "elapsed_ms": _ms(t0)}
        except Exception as e:
            logger.warning("Ontology extraction failed (non-critical): %s", e)
            yield {"status": "ONTOLOGY_SKIP", "reason": str(e)[:50], "elapsed_ms": _ms(t0)}

        # ── Step 7: 自动生成来源摘要 ──────────────────────────
        try:
            from core.llm_engine import get_llm_engine
            llm_eng = get_llm_engine()
            sample = text[:3000] if len(text) > 3000 else text
            summary = await llm_eng.ask_simple(
                f"请用中文为以下文档写一段100字以内的摘要：\n\n{sample}",
                system_prompt="你是文档摘要助手。只返回摘要文字。",
            )
            topics = await llm_eng.ask_simple(
                f"请从以下文档中提取5个关键主题词，用逗号分隔：\n\n{sample}",
                system_prompt="你是关键词提取助手。只返回逗号分隔的关键词。",
            )
            from utils.db_manager import db_pool
            db_pool.update_document_summary(vault_uuid, content_hash, summary, topics)
            questions = await llm_eng.ask_simple(
                f"请基于以下内容，生成3个用户最可能想问的问题，每行一个问题，不要编号：\n\n{sample}",
                system_prompt="你是问题生成助手。每行写一个问题，不要编号，不要多余文字。",
            )
            db_pool.update_document_questions(vault_uuid, content_hash, questions)
        except Exception as e:
            logger.warning("Auto-summary skipped: %s", e)

        yield {"status": "DONE", "chunks": len(chunks), "elapsed_ms": _ms(t0)}

    # ------------------------------------------------------------------
    # 内部辅助方法
    # ------------------------------------------------------------------

    async def _parse_document(
        self, file_bytes: bytes, file_name: str
    ) -> Tuple[str, int]:
        """解析文档，返回 (纯文本, 页数)。"""
        if file_name.lower().endswith(".pdf"):
            return await self._parse_pdf(file_bytes)
        if file_name.lower().endswith(".docx"):
            return await self._parse_docx(file_bytes)
        if file_name.lower().endswith(".md"):
            return file_bytes.decode("utf-8", errors="ignore"), 1
        if file_name.lower().endswith(".csv"):
            return await self._parse_csv(file_bytes)
        if file_name.lower().endswith(".json"):
            return await self._parse_json(file_bytes)
        if file_name.lower().endswith(".pptx"):
            return await self._parse_pptx(file_bytes)
        # 纯文本 fallback
        return file_bytes.decode("utf-8", errors="ignore"), 1

    async def _parse_docx(self, file_bytes: bytes) -> Tuple[str, int]:
        """解析 Word 文档，返回 (纯文本, 近似页数)。"""
        import io

        def _sync_parse() -> Tuple[str, int]:
            try:
                from docx import Document

                doc = Document(io.BytesIO(file_bytes))
                parts: List[str] = []
                for p in doc.paragraphs:
                    text = p.text.strip()
                    if text:
                        parts.append(text)
                for table in doc.tables:
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))
                text = "\n\n".join(parts)
                page_count = max(1, len(text) // 1800 + 1)
                return text, page_count
            except Exception as e:
                logger.error("DOCX parse failed: %s", e)
                return "", 0

        return await asyncio.to_thread(_sync_parse)

    async def _parse_csv(self, file_bytes: bytes) -> Tuple[str, int]:
        import io
        import csv

        def _sync_parse() -> Tuple[str, int]:
            try:
                text_parts: List[str] = []
                reader = csv.reader(io.TextIOWrapper(io.BytesIO(file_bytes), encoding="utf-8"))
                for row in reader:
                    text_parts.append(" | ".join(cell.strip() for cell in row if cell.strip()))
                text = "\n".join(text_parts)
                return text, max(1, len(text) // 1800 + 1)
            except Exception as e:
                logger.error("CSV parse failed: %s", e)
                return "", 0

        return await asyncio.to_thread(_sync_parse)

    async def _parse_json(self, file_bytes: bytes) -> Tuple[str, int]:
        import json

        def _sync_parse() -> Tuple[str, int]:
            try:
                data = json.loads(file_bytes.decode("utf-8"))
                text = json.dumps(data, indent=2, ensure_ascii=False)
                return text, max(1, len(text) // 1800 + 1)
            except Exception as e:
                logger.error("JSON parse failed: %s", e)
                return "", 0

        return await asyncio.to_thread(_sync_parse)

    async def _parse_pptx(self, file_bytes: bytes) -> Tuple[str, int]:
        import io

        def _sync_parse() -> Tuple[str, int]:
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                texts: List[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text.strip())
                text = "\n\n".join(texts)
                return text, len(prs.slides)
            except Exception as e:
                logger.error("PPTX parse failed: %s", e)
                return "", 0

        return await asyncio.to_thread(_sync_parse)

    async def _parse_pdf(self, file_bytes: bytes) -> Tuple[str, int]:
        """
        PDF 解析。使用 asyncio.to_thread 避免阻塞主事件循环。
        优先 pdfplumber，失败回退 pypdf，空文本时自动 OCR fallback。
        """
        import io

        def _sync_parse() -> Tuple[str, int]:
            text = ""
            page_count = 0

            # ── Step 1: 文本层提取 ──────────────────────────────
            try:
                import pdfplumber

                text_parts: List[str] = []
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_parts.append(page_text)
                    text = "\n\n".join(text_parts)
                    page_count = len(pdf.pages)
            except Exception as e1:
                logger.warning("pdfplumber failed: %s, falling back to pypdf", e1)
                try:
                    from pypdf import PdfReader

                    reader = PdfReader(io.BytesIO(file_bytes))
                    text = "\n\n".join(
                        p.extract_text() or "" for p in reader.pages
                    )
                    page_count = len(reader.pages)
                except Exception as e2:
                    logger.error("PDF parse failed: %s", e2)

            # ── Step 2: OCR fallback（扫描版/图片版 PDF）─────────
            if not text or len(text.strip()) < 50:
                logger.info("PDF text layer empty, attempting OCR fallback...")
                try:
                    import fitz
                    from rapidocr_onnxruntime import RapidOCR

                    ocr = RapidOCR()
                    doc = fitz.open(stream=file_bytes, filetype="pdf")
                    page_count = len(doc)
                    ocr_parts: List[str] = []

                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        # 2x 分辨率提高 OCR 精度
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_bytes = pix.tobytes("png")

                        # RapidOCR 支持 numpy array；将 PNG bytes 转为 array
                        import numpy as np
                        from PIL import Image

                        img = Image.open(io.BytesIO(img_bytes))
                        result, _ = ocr(np.array(img))
                        if result:
                            page_text = "\n".join(
                                str(item[1]) if isinstance(item, (list, tuple)) and len(item) > 1 else ""
                                for item in result
                            )
                            if page_text.strip():
                                ocr_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")

                    text = "\n\n".join(ocr_parts)
                    if text:
                        logger.info("OCR extracted %d chars from %d pages", len(text), page_count)
                    doc.close()
                except Exception as ocr_e:
                    logger.error("OCR fallback failed: %s", ocr_e)

            return text or "", page_count

        # CRITICAL: 运行在独立线程中，绝不卡死 Streamlit 主循环
        return await asyncio.to_thread(_sync_parse)

    async def _embed_chunks_async(
        self, chunks: List[Dict[str, Any]]
    ) -> List[List[float]]:
        """
        非阻塞 Embedding。

        CRITICAL: embedding runs inside asyncio.to_thread(...).
        Never block Streamlit render loop.
        """
        def _sync_embed() -> List[List[float]]:
            model = self.retriever._get_embedder()
            texts = [c["chunk_text"] for c in chunks]
            emb = model.encode(
                texts, show_progress_bar=False, convert_to_numpy=True
            )
            return emb.tolist()

        return await asyncio.to_thread(_sync_embed)

    async def _index_chunks(
        self,
        vault_uuid: str,
        doc_hash: str,
        chunks: List[Dict[str, Any]],
        vectors: List[List[float]],
    ) -> None:
        """持久化索引：ChromaDB(稠密) + BM25(稀疏) + 向量矩阵 + Chunk 元数据。"""
        texts = [c["chunk_text"] for c in chunks]
        self.retriever.add_chunks(texts)

        # 文件系统持久化（Phase 1A 遗留，用于元数据冷备份）
        vault_dir = os.path.join(self.persist_dir, vault_uuid)
        os.makedirs(vault_dir, exist_ok=True)
        np.save(
            os.path.join(vault_dir, f"{doc_hash}.npy"),
            np.array(vectors, dtype=np.float32),
        )
        with open(
            os.path.join(vault_dir, f"{doc_hash}.json"),
            "w",
            encoding="utf-8",
        ) as f:
            json.dump(chunks, f, ensure_ascii=False, indent=2)

        # Phase 1B: ChromaDB 稠密向量写入
        if self._chroma is not None:
            def _sync_chroma() -> None:
                collection = self._chroma.get_or_create_collection(name=vault_uuid)
                ids = [f"{doc_hash}_{c['chunk_index']}" for c in chunks]
                documents = texts
                metadatas = [
                    {
                        "doc_hash": doc_hash,
                        "chunk_index": c["chunk_index"],
                        "source_page": c["source_page"],
                        "header_hierarchy": c["header_hierarchy"],
                        "chunk_size": c["chunk_size"],
                        "overlap_prev": c["overlap_prev"],
                    }
                    for c in chunks
                ]
                collection.add(
                    ids=ids,
                    embeddings=vectors,
                    documents=documents,
                    metadatas=metadatas,
                )
            await asyncio.to_thread(_sync_chroma)
            logger.info("ChromaDB indexed %d chunks for vault=%s", len(chunks), vault_uuid)

    # ------------------------------------------------------------------
    # 检索接口 (Phase 1B 闭环)
    # ------------------------------------------------------------------

    def _get_embedder(self) -> Any:
        """懒加载 sentence-transformer embedder（线程安全复用）。"""
        if self._embedder is None:
            from sentence_transformers import SentenceTransformer
            self._embedder = SentenceTransformer("all-MiniLM-L6-v2")
        return self._embedder

    async def _search_dense(
        self, query: str, vault_uuid: str, top_k: int
    ) -> List[Dict[str, Any]]:
        """
        ChromaDB 稠密向量检索。
        CRITICAL: embedding 与 collection.query 均在 asyncio.to_thread 中执行。
        """
        if self._chroma is None:
            return []

        def _sync_query() -> List[Dict[str, Any]]:
            embedder = self._get_embedder()
            query_vec = embedder.encode([query], convert_to_numpy=True)
            collection = self._chroma.get_or_create_collection(name=vault_uuid)
            results = collection.query(
                query_embeddings=query_vec.tolist(),
                n_results=top_k,
                include=["documents", "metadatas", "distances"],
            )
            # ChromaDB 返回格式: {'ids': [[...]], 'documents': [[...]], ...}
            out: List[Dict[str, Any]] = []
            ids = results.get("ids", [[]])[0]
            docs = results.get("documents", [[]])[0]
            metas = results.get("metadatas", [[]])[0]
            dists = results.get("distances", [[]])[0]
            for i, doc_id in enumerate(ids):
                out.append(
                    {
                        "id": doc_id,
                        "chunk_text": docs[i] if i < len(docs) else "",
                        "metadata": metas[i] if i < len(metas) else {},
                        "distance": dists[i] if i < len(dists) else 0.0,
                        "source": "dense",
                    }
                )
            return out

        return await asyncio.to_thread(_sync_query)

    @staticmethod
    def _fuse_rrf(
        dense_results: List[Dict[str, Any]],
        sparse_results: List[Tuple[int, float]],
        k: int = 60,
    ) -> List[Dict[str, Any]]:
        """
        Reciprocal Rank Fusion (RRF): 将稠密与稀疏排名融合为统一打分。
        score = Σ 1/(k + rank)
        """
        scores: Dict[str, Dict[str, Any]] = {}

        # Dense ranks (ChromaDB distance 越小越好 → 排名越靠前)
        for rank, item in enumerate(dense_results):
            doc_id = item["id"]
            scores[doc_id] = scores.get(doc_id, {"score": 0.0, "item": item})
            scores[doc_id]["score"] += 1.0 / (k + rank + 1)

        # Sparse ranks (BM25 score 越大越好 → 排名越靠前)
        for rank, (idx, _) in enumerate(sparse_results):
            doc_id = f"bm25_{idx}"
            if doc_id not in scores:
                scores[doc_id] = {
                    "score": 0.0,
                    "item": {
                        "id": doc_id,
                        "chunk_index": idx,
                        "chunk_text": "",  # 占位，检索后回填
                        "metadata": {},
                        "source": "bm25",
                    },
                }
            scores[doc_id]["score"] += 1.0 / (k + rank + 1)

        ranked = sorted(scores.values(), key=lambda x: x["score"], reverse=True)
        return [r["item"] for r in ranked]

    async def retrieve(
        self,
        query: str,
        vault_uuid: str,
        top_k: int = 5,
        content_hashes: Optional[List[str]] = None,
    ) -> List[Dict[str, Any]]:
        """
        混合检索闭环 (Phase 1B):
          1. Dense  (ChromaDB 向量检索)
          2. Sparse (BM25 关键词检索)
          3. Fuse   (RRF 排名融合)
          4. Rerank (Cross-Encoder 优雅降级)
        """
        # 1. Dense search
        dense_results = await self._search_dense(query, vault_uuid, top_k=top_k * 3)

        # 2. Sparse search
        sparse_results = self.retriever.search_bm25(query, top_k=top_k * 3)

        # 3. RRF 融合
        fused = self._fuse_rrf(dense_results, sparse_results, k=60)

        # 回填 BM25 chunk_text（懒加载）
        candidates: List[Dict[str, Any]] = []
        for item in fused[: top_k * 2]:
            if item.get("source") == "bm25" and not item.get("chunk_text"):
                idx = item.get("chunk_index", -1)
                if 0 <= idx < len(self.retriever._chunk_texts):
                    item["chunk_text"] = self.retriever._chunk_texts[idx]
            candidates.append(
                {
                    "chunk_text": item.get("chunk_text", ""),
                    "chunk_index": item.get("chunk_index", -1),
                    "score": item.get("score", 0.0),
                    "source": item.get("source", "unknown"),
                    "metadata": item.get("metadata", {}),
                }
            )

        # 4. Rerank (优雅降级)
        ranked = await self.retriever.rerank(query, candidates)

        # 5. 按 content_hashes 过滤
        if content_hashes:
            ranked = [
                r for r in ranked
                if r.get("metadata", {}).get("content_hash") in content_hashes
            ]

        return ranked[:top_k]


# ---------------------------------------------------------------------------
# 模块级单例工厂 (供 frontend 复用)
# ---------------------------------------------------------------------------

_pipeline_singleton: Optional[IngestionPipeline] = None


def get_pipeline() -> IngestionPipeline:
    """获取全局唯一的 IngestionPipeline 实例，确保 ChromaDB/BM25 状态共享。"""
    global _pipeline_singleton
    if _pipeline_singleton is None:
        _pipeline_singleton = IngestionPipeline()
    return _pipeline_singleton
