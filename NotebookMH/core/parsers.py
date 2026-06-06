"""core/parsers.py — 文件/URL 解析"""
import csv, json, io, logging, re
from typing import Optional

log = logging.getLogger(__name__)


def parse_pdf(data: bytes) -> dict:
    import pdfplumber
    page_count = 0
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            page_count = len(pdf.pages)
            pages = [p.extract_text() or "" for p in pdf.pages]
        text = "\n\n".join(pages)
        if text.strip():
            return {"text": text, "page_count": page_count}
        # 无文本层，尝试 OCR（扫描件/图片型 PDF）
        try:
            from pdf2image import convert_from_bytes
            import pytesseract
            images = convert_from_bytes(
                data, dpi=200,
                first_page=1, last_page=min(page_count, 5)
            )
            ocr_parts = [pytesseract.image_to_string(img, lang="chi_sim+eng")
                         for img in images]
            ocr_text = "\n\n".join(ocr_parts)
            if ocr_text.strip():
                return {"text": ocr_text, "page_count": len(images), "ocr": True}
        except ImportError:
            pass
        except Exception as exc:
            return {"text": "", "page_count": page_count,
                    "error": f"扫描件 OCR 失败: {exc}。"
                             f"请安装依赖: pip install pdf2image pytesseract pillow，"
                             f"并安装 Tesseract-OCR 引擎。"}
        return {"text": "", "page_count": page_count,
                "error": "PDF 无文本层（可能是扫描件/图片型 PDF）。"
                         "若需识别扫描件，请安装 OCR 依赖（见 README）。"}
    except Exception as exc:
        return {"text": "", "page_count": page_count, "error": str(exc)}


def parse_docx(data: bytes) -> dict:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return {"text": "\n".join(paras), "page_count": 0}


def parse_pptx(data: bytes) -> dict:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        slides.append("\n".join(texts))
    return {"text": "\n\n".join(slides), "page_count": len(slides)}


def parse_txt(data: bytes) -> dict:
    for enc in ("utf-8", "gbk", "gb18030", "latin-1"):
        try:
            return {"text": data.decode(enc), "page_count": 0}
        except UnicodeDecodeError:
            continue
    return {"text": data.decode("utf-8", errors="ignore"), "page_count": 0}


def parse_md(data: bytes) -> dict:
    return parse_txt(data)


def parse_csv(data: bytes) -> dict:
    text = parse_txt(data)["text"]
    reader = csv.reader(io.StringIO(text))
    rows = ["\t".join(r) for r in reader]
    return {"text": "\n".join(rows), "page_count": 0}


def parse_json(data: bytes) -> dict:
    text = parse_txt(data)["text"]
    try:
        obj = json.loads(text)
        return {"text": json.dumps(obj, ensure_ascii=False, indent=2), "page_count": 0}
    except json.JSONDecodeError:
        return {"text": text, "page_count": 0}


_URL_HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/120.0.0.0 Safari/537.36 Edg/120.0.0.0"
    ),
    "Accept": (
        "text/html,application/xhtml+xml,application/xml;"
        "q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8"
    ),
    "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8,en-GB;q=0.7,en-US;q=0.6",
    "Referer": "https://www.google.com/",
    "Connection": "keep-alive",
    "Upgrade-Insecure-Requests": "1",
}

# 静态抓取文本低于此阈值时，认为可能是 JS 渲染页，触发无头浏览器兑底
_MIN_STATIC_TEXT = 200


def _extract_main_text(html: str) -> str:
    """从 HTML 中提取正文，优先主体容器，去除导航/广告/脚本等噪音。"""
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header",
                     "aside", "form", "iframe", "noscript", "svg",
                     "button", "input"]):
        tag.decompose()
    # 优先从语义主体容器提取
    main = (soup.find("article") or soup.find("main")
            or soup.find(attrs={"role": "main"}))
    target = main or soup.body or soup
    text = target.get_text(separator="\n", strip=True)
    # 压缩多余空行
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return re.sub(r"\n{3,}", "\n\n", "\n".join(lines))


def _fetch_static(url: str):
    """httpx 抓静态 HTML。返回 (html, error)。"""
    import httpx
    try:
        with httpx.Client(timeout=10, follow_redirects=True,
                          headers=_URL_HEADERS) as client:
            r = client.get(url)
            r.raise_for_status()
            return r.text, None
    except httpx.HTTPStatusError as e:
        return "", f"HTTP {e.response.status_code}，无法访问该网页"
    except httpx.RequestError as e:
        return "", f"网络请求失败: {e}"
    except Exception as e:
        return "", f"网页获取失败: {e}"


def _fetch_headless(url: str) -> str:
    """Playwright 无头浏览器渲染 JS 后取 HTML。未安装则抛 ImportError。"""
    from playwright.sync_api import sync_playwright
    with sync_playwright() as p:
        browser = p.chromium.launch(
            headless=True,
            args=["--no-sandbox", "--disable-dev-shm-usage"],
        )
        try:
            page = browser.new_page(user_agent=_URL_HEADERS["User-Agent"])
            page.goto(url, wait_until="networkidle", timeout=10000)
            html = page.content()
        finally:
            browser.close()
    return html


def parse_url(url: str) -> dict:
    html, err = _fetch_static(url)
    text = _extract_main_text(html) if html else ""

    # 静态抓取太少或报错 → 尝试无头浏览器兑底（JS 渲染站点）
    if len(text) < _MIN_STATIC_TEXT:
        try:
            html2 = _fetch_headless(url)
            text2 = _extract_main_text(html2)
            if len(text2) > len(text):
                text, err = text2, None
        except ImportError:
            log.info("Playwright 未安装，跳过无头浏览器兑底")
        except Exception as e:
            log.warning("无头浏览器抓取失败: %s", e)

    if text.strip():
        return {"text": text, "page_count": 0}
    if err:
        return {"text": "", "page_count": 0, "error": err}
    return {"text": "", "page_count": 0,
            "error": "该网页可能是 JS 动态渲染或需登录，未能提取到正文"}


_DISPATCH = {
    "pdf": parse_pdf, "docx": parse_docx, "pptx": parse_pptx,
    "txt": parse_txt, "md": parse_md, "csv": parse_csv, "json": parse_json,
}


def parse_file(file_name: str, data: bytes) -> dict:
    ext = file_name.rsplit(".", 1)[-1].lower()
    if ext not in _DISPATCH:
        raise ValueError(f"不支持的文件类型: {ext}")
    return _DISPATCH[ext](data)
