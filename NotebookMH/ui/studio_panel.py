"""ui/studio_panel.py — Studio 右侧面板"""
import asyncio
import traceback

import streamlit as st

from core.db import db_manager

# 工具元数据：标题 → studio 模块函数名
_TOOLS = [
    ("summary",       "📝 摘要",       "generate_summary"),
    ("faq",           "❓ 常见问题",   "generate_faq"),
    ("study_guide",   "📚 学习指南",   "generate_study_guide"),
    ("briefing",      "📰 简报",       "generate_briefing"),
    ("timeline",      "🕒 时间线",     "generate_timeline"),
    ("mindmap",       "🧠 思维导图",   "generate_mindmap"),
    ("presentation",  "📊 演示文稿",   "generate_presentation"),
    ("flashcards",    "🃏 闪卡",       "generate_flashcards"),
    ("quiz",          "📋 测验",       "generate_quiz"),
]


def _run_tool(tool_key: str, vault_uuid: str):
    from core import studio as studio_mod
    func_name = dict((k, fn) for k, _, fn in _TOOLS)[tool_key]
    fn = getattr(studio_mod, func_name)
    return asyncio.run(fn(vault_uuid))


def _save_note(vault_uuid: str, user_id: str, title: str, content: str):
    db_manager.save_note(vault_uuid, user_id, title, content)


def _render_tools_grid(vault_uuid: str, user_id: str) -> None:
    cols = st.columns(3)
    for i, (key, label, _) in enumerate(_TOOLS):
        with cols[i % 3]:
            if st.button(label, key=f"studio_tool_{key}",
                         use_container_width=True):
                st.session_state[f"_studio_running_{key}"] = True
                st.rerun()

    # 真正的执行（避免 button 后 rerun 丢失结果）
    for key, label, _ in _TOOLS:
        running_flag = f"_studio_running_{key}"
        if st.session_state.get(running_flag):
            st.session_state[running_flag] = False
            # 单选模式：清除其他工具的结果，只保留当前
            for other_key, _, _ in _TOOLS:
                if other_key != key:
                    st.session_state.pop(f"_studio_result_{other_key}", None)
            with st.spinner(f"生成 {label}..."):
                try:
                    result = _run_tool(key, vault_uuid)
                    st.session_state[f"_studio_result_{key}"] = result
                    # 新生成的逐题内容，重置交互进度
                    if key == "flashcards":
                        for k in ("_fc_idx", "_fc_show", "_fc_known"):
                            st.session_state.pop(k, None)
                    elif key == "quiz":
                        for k in ("_qz_idx", "_qz_answered", "_qz_score",
                                  "_qz_last_correct"):
                            st.session_state.pop(k, None)
                except Exception:
                    st.session_state[f"_studio_result_{key}"] = traceback.format_exc()
            st.rerun()

    # 渲染结果
    for key, label, _ in _TOOLS:
        result_key = f"_studio_result_{key}"
        if result_key in st.session_state:
            res = st.session_state[result_key]
            with st.expander(f"{label} 结果", expanded=True):
                _render_result(key, label, res, vault_uuid, user_id)


def _render_result(key: str, label: str, res, vault_uuid: str, user_id: str):
    if isinstance(res, str) and res.startswith("Traceback"):
        st.error("生成失败：")
        st.code(res)
        return

    if key == "mindmap":
        import html
        safe = html.escape(res)
        st.markdown("**Mermaid 源码:**")
        with st.expander("查看源码"):
            st.code(res, language="mermaid")
        st.components.v1.html(
            f"""
<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<style>
  body {{ margin: 0; background: #f8fafc; }}
  #mindmap-wrap {{ width: 100%; height: 100%; overflow: hidden; position: relative; }}
  #mindmap-svg {{ width: 100%; height: 100%; }}
</style>
</head>
<body>
<div id="mindmap-wrap">
  <div class="mermaid">{safe}</div>
</div>
<script src="https://cdn.jsdelivr.net/npm/mermaid@10/dist/mermaid.min.js"></script>
<script>
  mermaid.initialize({{
    startOnLoad: true,
    theme: 'base',
    themeVariables: {{
      primaryColor: '#dbeafe',
      primaryTextColor: '#1e3a5f',
      primaryBorderColor: '#3b82f6',
      lineColor: '#64748b',
      secondaryColor: '#f1f5f9',
      tertiaryColor: '#ffffff',
      fontFamily: '"Inter", "Noto Sans SC", "PingFang SC", sans-serif',
      fontSize: '14px'
    }},
    flowchart: {{ useMaxWidth: true, htmlLabels: true, curve: 'basis' }},
    mindmap: {{ useMaxWidth: true }}
  }});
</script>
</body>
</html>
            """,
            height=600, scrolling=True,
        )
        cols = st.columns(2)
        if cols[0].button("保存为笔记", key=f"save_{key}"):
            _save_note(vault_uuid, user_id, label, f"```mermaid\n{res}\n```")
            st.success("已保存")
        if cols[1].button("清除", key=f"clear_{key}"):
            del st.session_state[f"_studio_result_{key}"]
            st.rerun()

    elif key == "flashcards":
        if not isinstance(res, list) or not res:
            st.warning("未生成闪卡")
            return
        _render_flashcard_session(res, vault_uuid)

    elif key == "quiz":
        if not isinstance(res, list) or not res:
            st.warning("未生成测验题")
            return
        _render_quiz_session(res, vault_uuid)

    elif key == "presentation":
        slides = res.get("slides") if isinstance(res, dict) else []
        if not slides:
            st.warning("未生成演示文稿")
            return
        st.markdown(f"共 {len(slides)} 页")
        for i, s in enumerate(slides):
            with st.container(border=True):
                st.markdown(f"**{i+1}. {s['title']}**")
                for b in s.get("bullets", []):
                    st.markdown(f"- {b}")
                notes = s.get("speaker_notes")
                if notes:
                    st.caption(f"备注: {notes}")
        pptx_bytes = _to_pptx(slides)
        cols = st.columns(3)
        cols[0].download_button(
            "下载 PPT", pptx_bytes,
            file_name="presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"dl_pptx_{key}")
        if cols[1].button("保存为笔记", key=f"save_{key}"):
            text = "\n\n".join(
                f"**{s['title']}**\n" + "\n".join(f"- {b}" for b in s.get("bullets", []))
                for s in slides
            )
            _save_note(vault_uuid, user_id, label, text)
            st.success("已保存")
        if cols[2].button("清除", key=f"clear_{key}"):
            del st.session_state[f"_studio_result_{key}"]
            st.rerun()

    else:
        # 纯文本（摘要 / FAQ / 学习指南 / 简报 / 时间线）
        st.markdown(res)
        cols = st.columns(2)
        if cols[0].button("保存为笔记", key=f"save_{key}"):
            _save_note(vault_uuid, user_id, label, res)
            st.success("已保存")
        if cols[1].button("清除", key=f"clear_{key}"):
            del st.session_state[f"_studio_result_{key}"]
            st.rerun()


def _render_flashcard_session(cards: list, vault_uuid: str) -> None:
    """闪卡逐张交互：显示问题 → 看答案 → 记住了/没记住 → 下一张。"""
    st.session_state.setdefault("_fc_idx", 0)
    st.session_state.setdefault("_fc_show", False)
    st.session_state.setdefault("_fc_known", 0)
    n = len(cards)
    idx = st.session_state["_fc_idx"]

    if idx >= n:
        known = st.session_state["_fc_known"]
        st.success(f"本轮完成！共 {n} 张，记住 {known} 张")
        c0, c1, c2 = st.columns(3)
        if c0.button("🔄 再来一轮", key="fc_restart", use_container_width=True):
            st.session_state["_fc_idx"] = 0
            st.session_state["_fc_show"] = False
            st.session_state["_fc_known"] = 0
            st.rerun()
        if c1.button("💾 保存到闪卡库", key="fc_save_lib",
                     use_container_width=True):
            db_manager.save_flashcards(vault_uuid, cards)
            st.success(f"已保存 {n} 张")
        if c2.button("清除", key="fc_clear", use_container_width=True):
            for k in ("_fc_idx", "_fc_show", "_fc_known"):
                st.session_state.pop(k, None)
            st.session_state.pop("_studio_result_flashcards", None)
            st.rerun()
        return

    card = cards[idx]
    st.progress((idx) / n, text=f"第 {idx + 1} / {n} 张")
    with st.container(border=True):
        st.markdown(f"#### {card['question']}")
        if not st.session_state["_fc_show"]:
            if st.button("👀 显示答案", key=f"fc_show_{idx}",
                         use_container_width=True):
                st.session_state["_fc_show"] = True
                st.rerun()
        else:
            st.info(card["answer"])
            c0, c1 = st.columns(2)
            if c0.button("✅ 记住了", key=f"fc_known_{idx}",
                         use_container_width=True):
                st.session_state["_fc_known"] += 1
                st.session_state["_fc_idx"] += 1
                st.session_state["_fc_show"] = False
                st.rerun()
            if c1.button("❌ 没记住", key=f"fc_again_{idx}",
                         use_container_width=True):
                st.session_state["_fc_idx"] += 1
                st.session_state["_fc_show"] = False
                st.rerun()


def _render_quiz_session(items: list, vault_uuid: str) -> None:
    """测验交互：支持选择题/填空题/读图题/综合题，可逐题作答或查看完整试卷。"""
    st.session_state.setdefault("_qz_idx", 0)
    st.session_state.setdefault("_qz_answered", False)
    st.session_state.setdefault("_qz_score", 0)
    st.session_state.setdefault("_qz_view_mode", "interactive")  # interactive | full_paper
    n = len(items)

    # ── 模式切换 ──
    mode = st.session_state["_qz_view_mode"]
    c0, c1 = st.columns(2)
    if c0.button("📋 逐题作答模式", key="qz_mode_interactive",
                 use_container_width=True, type="primary" if mode == "interactive" else "secondary"):
        st.session_state["_qz_view_mode"] = "interactive"
        st.rerun()
    if c1.button("📄 查看完整试卷", key="qz_mode_full",
                 use_container_width=True, type="primary" if mode == "full_paper" else "secondary"):
        st.session_state["_qz_view_mode"] = "full_paper"
        st.rerun()

    if mode == "full_paper":
        _render_full_paper(items)
        return

    # ── 逐题作答模式 ──
    idx = st.session_state["_qz_idx"]
    if idx >= n:
        score = st.session_state["_qz_score"]
        total_score = sum(it.get("score", 1) for it in items)
        st.success(f"测验完成！得分 {score} / {total_score}")
        c0, c1 = st.columns(2)
        if c0.button("🔄 再来一轮", key="qz_restart", use_container_width=True):
            st.session_state["_qz_idx"] = 0
            st.session_state["_qz_answered"] = False
            st.session_state["_qz_score"] = 0
            st.rerun()
        if c1.button("清除", key="qz_clear", use_container_width=True):
            for k in ("_qz_idx", "_qz_answered", "_qz_score",
                      "_qz_last_correct", "_qz_view_mode"):
                st.session_state.pop(k, None)
            st.session_state.pop("_studio_result_quiz", None)
            st.rerun()
        return

    it = items[idx]
    q_type = it.get("type", "选择")
    q_score = it.get("score", 1)
    st.progress(idx / n, text=f"第 {idx + 1} / {n} 题  |  {q_type}（{q_score}分）")
    with st.container(border=True):
        st.markdown(f"**【{q_type}】** {it['question']}")
        _render_visual(it.get("visual"), key=f"vis_q_{idx}")
        answered = st.session_state["_qz_answered"]
        user_answer = None

        if "选择" in q_type:
            opts = it.get("options", [])
            if opts:
                user_answer = st.radio("选择答案", opts, key=f"qz_opt_{idx}",
                                       index=None, disabled=answered)
            else:
                st.info("（本题无选项，请直接作答）")
                user_answer = st.text_input("你的答案", key=f"qz_text_{idx}",
                                           disabled=answered)
        else:
            user_answer = st.text_area("你的答案", key=f"qz_text_{idx}",
                                       height=80, disabled=answered)

        if not answered:
            if st.button("提交", key=f"qz_submit_{idx}",
                         use_container_width=True):
                if not user_answer or (isinstance(user_answer, str) and not user_answer.strip()):
                    st.warning("请先作答")
                else:
                    correct = str(it["correct"]).strip()
                    ans_str = user_answer.strip() if isinstance(user_answer, str) else str(user_answer).strip()
                    # 宽松匹配：选择题取首字母，其他题去除空格后比对
                    if "选择" in q_type and len(ans_str) == 1:
                        ans_clean = ans_str.upper()
                        corr_clean = correct.upper()
                    else:
                        ans_clean = ans_str.replace(" ", "").replace("，", ",").lower()
                        corr_clean = correct.replace(" ", "").replace("，", ",").lower()
                    ok = ans_clean == corr_clean
                    st.session_state["_qz_answered"] = True
                    st.session_state["_qz_last_correct"] = ok
                    if ok:
                        st.session_state["_qz_score"] += q_score
                    st.rerun()
        else:
            if st.session_state.get("_qz_last_correct"):
                st.success("✅ 答对了！")
            else:
                st.error(f"❌ 答错了，正确答案：{it['correct']}")
            if it.get("explanation"):
                st.info(f"**解析：** {it['explanation']}")
            if st.button("下一题 ▶", key=f"qz_next_{idx}",
                         use_container_width=True):
                st.session_state["_qz_idx"] += 1
                st.session_state["_qz_answered"] = False
                st.rerun()


def _render_full_paper(items: list) -> None:
    """展示完整试卷（所有题目+答案+解析）。"""
    st.markdown("### 📄 完整模拟试卷")
    total = sum(it.get("score", 1) for it in items)
    st.caption(f"共 {len(items)} 题，满分 {total} 分")

    for i, it in enumerate(items, 1):
        q_type = it.get("type", "选择")
        score = it.get("score", 1)
        with st.container(border=True):
            st.markdown(f"**{i}. 【{q_type} · {score}分】** {it['question']}")
            _render_visual(it.get("visual"), key=f"vis_p_{i}")
            if "选择" in q_type and it.get("options"):
                for opt in it["options"]:
                    st.markdown(f"&nbsp;&nbsp;&nbsp;&nbsp;{opt}")
            with st.expander("查看答案与解析"):
                st.markdown(f"**答案：** {it['correct']}")
                if it.get("explanation"):
                    st.markdown(f"**解析：** {it['explanation']}")


def _render_visual(visual, key: str = "") -> None:
    """渲染题目配图：SVG 直接嵌入，image 显示网络图片。"""
    if not isinstance(visual, dict):
        return
    kind = visual.get("kind")
    if kind == "svg" and visual.get("svg"):
        st.components.v1.html(
            f'<div style="background:#fff;padding:8px;border-radius:8px;">{visual["svg"]}</div>',
            height=340, scrolling=True,
        )
    elif kind == "image" and visual.get("url"):
        try:
            st.image(visual["url"], use_container_width=True)
        except Exception:
            st.caption(f"配图加载失败：{visual['url']}")


def _render_quiz_library(vault_uuid: str) -> None:
    items = db_manager.list_quiz_items(vault_uuid, only_unanswered=True)
    if not items:
        return
    with st.expander(f"待答测验 ({len(items)} 题)"):
        for q in items[:10]:
            with st.container(border=True):
                st.markdown(f"**Q:** {q.question}")
                for opt in (q.options or []):
                    st.markdown(f"- {opt}")
                ans = st.radio("你的答案", options=["A","B","C","D"],
                               key=f"qz_{q.id}", horizontal=True,
                               label_visibility="collapsed")
                if st.button("提交", key=f"qz_submit_{q.id}"):
                    ok = db_manager.answer_quiz(q.id, ans)
                    if ok:
                        st.success("答对了！")
                    else:
                        st.error(f"答错了。正确答案: {q.correct}")
                        st.caption(q.explanation or "")
                    st.rerun()


def _render_wrong_answers(vault_uuid: str) -> None:
    wrongs = db_manager.list_wrong_answers(vault_uuid, only_unmastered=True)
    if not wrongs:
        return
    with st.expander(f"错题本 ({len(wrongs)} 题)"):
        for w in wrongs[:20]:
            with st.container(border=True):
                st.markdown(f"**Q:** {w.question}")
                st.caption(f"你的答案: {w.user_answer} | 正确: {w.correct_answer}")
                if w.explanation:
                    st.markdown(f"**解析**: {w.explanation}")
                if st.button("已掌握", key=f"wa_{w.id}"):
                    db_manager.mark_wrong_mastered(w.id)
                    st.rerun()


def _render_flashcard_library(vault_uuid: str) -> None:
    cards = db_manager.list_flashcards(vault_uuid)
    if not cards:
        return
    with st.expander(f"闪卡库 ({len(cards)} 张)"):
        for c in cards[:20]:
            with st.container(border=True):
                st.markdown(f"**Q:** {c.question}")
                with st.expander("答案"):
                    st.markdown(c.answer)
                    cols = st.columns(3)
                    if cols[0].button("未掌握", key=f"fc0_{c.id}"):
                        db_manager.update_flashcard_mastery(c.id, 0); st.rerun()
                    if cols[1].button("半懂", key=f"fc1_{c.id}"):
                        db_manager.update_flashcard_mastery(c.id, 1); st.rerun()
                    if cols[2].button("已掌握", key=f"fc2_{c.id}"):
                        db_manager.update_flashcard_mastery(c.id, 2); st.rerun()


def _to_pptx(slides: list[dict]) -> bytes:
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    for s in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = s.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(s.get("bullets", [])):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _to_docx(title: str, content: str) -> bytes:
    import io
    from docx import Document
    doc = Document()
    doc.add_heading(title, level=1)
    for para in content.split("\n\n"):
        doc.add_paragraph(para)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _to_pdf(title: str, content: str) -> bytes:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        import io
        try:
            pdfmetrics.registerFont(TTFont("CJK", "C:/Windows/Fonts/simhei.ttf"))
            font_name = "CJK"
        except Exception:
            font_name = "Helvetica"
        buf = io.BytesIO()
        c = canvas.Canvas(buf, pagesize=A4)
        c.setFont(font_name, 16)
        c.drawString(50, 800, title)
        c.setFont(font_name, 11)
        y = 770
        for line in content.split("\n"):
            if y < 50:
                c.showPage()
                y = 800
                c.setFont(font_name, 11)
            c.drawString(50, y, line[:80])
            y -= 18
        c.save()
        return buf.getvalue()
    except ImportError:
        return f"{title}\n\n{content}".encode("utf-8")


def _render_notes_section(vault_uuid: str, user_id: str) -> None:
    st.markdown("### 我的笔记")
    notes = db_manager.list_notes(vault_uuid, user_id)
    if not notes:
        st.caption("还没有保存的笔记")
        return
    for n in notes[:20]:
        with st.expander(f"{'📌 ' if n.pinned else ''}{n.title}"):
            st.markdown(n.content)
            md_bytes = f"# {n.title}\n\n{n.content}".encode("utf-8")
            cols = st.columns(5)
            cols[0].download_button("MD", md_bytes, file_name=f"{n.title}.md",
                                    mime="text/markdown", key=f"dl_md_{n.id}")
            cols[1].download_button("Word", _to_docx(n.title, n.content),
                                    file_name=f"{n.title}.docx",
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                    key=f"dl_docx_{n.id}")
            cols[2].download_button("PDF", _to_pdf(n.title, n.content),
                                    file_name=f"{n.title}.pdf",
                                    mime="application/pdf",
                                    key=f"dl_pdf_{n.id}")
            if cols[3].button("置顶", key=f"pin_{n.id}"):
                db_manager.toggle_pin_note(n.id); st.rerun()
            if cols[4].button("删除", key=f"del_note_{n.id}"):
                db_manager.delete_note(n.id); st.rerun()


def render() -> None:
    st.markdown("### Studio")
    vault_uuid = st.session_state.get("vault_uuid", "")
    user_id = st.session_state.get("user_id", "anonymous")
    if not vault_uuid:
        st.caption("请先选择笔记库")
        return

    docs = db_manager.list_documents(vault_uuid)
    if not docs:
        st.caption("先上传资料再生成 Studio 内容")
        return

    _render_tools_grid(vault_uuid, user_id)
    st.divider()
    _render_quiz_library(vault_uuid)
    _render_wrong_answers(vault_uuid)
    _render_flashcard_library(vault_uuid)
    _render_notes_section(vault_uuid, user_id)
